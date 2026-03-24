from pptx import Presentation

prs = Presentation('MPL cricket.pptx')

print("Updated player list from PPT (46 slides):\n")
print("Serial | Name           | Role")
print("-" * 50)

players_list = []
for slide_idx, slide in enumerate(prs.slides):
    serial = slide_idx + 1  # Slide index + 1 = serial number
    name = ""
    role = ""
    
    # Find name and role in shapes
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text and text not in ["", "ALL-ROUNDER", "BATSMAN", "BOWLER", "WICKETKEEPER"]:
                if not name:
                    name = text
            elif text in ["ALL-ROUNDER", "BATSMAN", "BOWLER", "WICKETKEEPER"]:
                role = text
    
    if name and role:
        print(f"{serial:3d}    | {name:14s} | {role}")
        players_list.append((serial, name, role))

print(f"\nTotal: {len(players_list)} players")
