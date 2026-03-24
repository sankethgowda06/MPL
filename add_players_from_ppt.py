import os
from pptx import Presentation
from app import db, Player
import json

# Path to the PowerPoint file
ppt_path = r'c:\Users\Sanketh\Desktop\MPL\app\MPL cricket.pptx'
images_dir = r'c:\Users\Sanketh\Desktop\MPL\app\player_photos'

# Create images directory if it doesn't exist
os.makedirs(images_dir, exist_ok=True)

# Load the PowerPoint presentation
prs = Presentation(ppt_path)

# Extract player data
image_count = 0

# Extract player details and save for verification
print("Starting player extraction from PowerPoint...")
extracted_players = []

for slide_idx, slide in enumerate(prs.slides):
    print(f"\nProcessing Slide {slide_idx + 1}...")
    
    # Extract text from shapes
    texts = []
    image_path = None
    
    for shape_idx, shape in enumerate(slide.shapes):
        # Extract text
        if hasattr(shape, 'text') and shape.text.strip():
            text = shape.text.strip()
            texts.append(text)
            print(f"  Text found: {text}")
        
        # Extract images
        if shape.shape_type == 13:  # Picture shape type
            try:
                image = shape.image
                image_bytes = image.blob
                image_filename = f'player_{slide_idx + 1}.{image.ext}'
                image_filepath = os.path.join(images_dir, image_filename)
                
                with open(image_filepath, 'wb') as f:
                    f.write(image_bytes)
                
                image_path = image_filepath
                print(f"  Image extracted: {image_filename}")
                image_count += 1
            except Exception as e:
                print(f"  Error extracting image: {e}")
    
    # Parse texts to identify player name and role
    player_name = None
    role = None
    
    # Look for role keywords
    role_keywords = ['all-rounder', 'allrounder', 'batsman', 'bowler', 'keeper', 'wicket keeper', 'spinner']
    
    for text in texts:
        text_lower = text.lower()
        # Check if this text is a role
        is_role = any(keyword in text_lower for keyword in role_keywords)
        
        if is_role and role is None:
            role = text
        elif not is_role and player_name is None:
            player_name = text
    
    # Save player details for verification
    if player_name:
        extracted_players.append({
            'Serial Number': slide_idx + 1,
            'Player Name': player_name,
            'Role': role or 'Unknown',
            'Photo Path': image_path
        })
        print(f"  Player extracted: {player_name} - {role or 'Unknown'}")

# Save extracted players to a file for verification
output_file = os.path.join(images_dir, 'extracted_players.json')
with open(output_file, 'w') as f:
    json.dump(extracted_players, f, indent=4)

print(f"\n✓ Player extraction complete. Total players extracted: {len(extracted_players)}")
print(f"✓ Extracted player details saved to: {output_file}")

# Debugging player addition
print("Starting player addition from PowerPoint...")
with db.app.app_context():
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\nProcessing Slide {slide_idx + 1}...")
        
        # Extract text from shapes
        texts = []
        image_path = None
        
        for shape_idx, shape in enumerate(slide.shapes):
            # Extract text
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                texts.append(text)
                print(f"  Text found: {text}")
            
            # Extract images
            if shape.shape_type == 13:  # Picture shape type
                try:
                    image = shape.image
                    image_bytes = image.blob
                    image_filename = f'player_{slide_idx + 1}.{image.ext}'
                    image_filepath = os.path.join(images_dir, image_filename)
                    
                    with open(image_filepath, 'wb') as f:
                        f.write(image_bytes)
                    
                    image_path = image_filepath
                    print(f"  Image extracted: {image_filename}")
                    image_count += 1
                except Exception as e:
                    print(f"  Error extracting image: {e}")
        
        # Parse texts to identify player name and role
        player_name = None
        role = None
        
        # Look for role keywords
        role_keywords = ['all-rounder', 'allrounder', 'batsman', 'bowler', 'keeper', 'wicket keeper', 'spinner']
        
        for text in texts:
            text_lower = text.lower()
            # Check if this text is a role
            is_role = any(keyword in text_lower for keyword in role_keywords)
            
            if is_role and role is None:
                role = text
            elif not is_role and player_name is None:
                player_name = text
        
        # Add player to database if we have at least a name
        if player_name:
            print(f"  Attempting to add player: {player_name} - {role or 'Unknown'}")
            if not Player.query.filter_by(name=player_name).first():
                player = Player(
                    serial_number=slide_idx + 1,
                    name=player_name,
                    role=role or 'Unknown',
                    photo_path=image_path,
                    is_available=True
                )
                db.session.add(player)
                print(f"  Player added: {player_name} - {role or 'Unknown'}")
            else:
                print(f"  Player already exists: {player_name}")
    
    db.session.commit()
    print(f"\n✓ All players added successfully. Total images extracted: {image_count}")
