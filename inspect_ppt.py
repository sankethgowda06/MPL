from pptx import Presentation

prs = Presentation('MPL cricket.pptx')

print(f"Total slides: {len(prs.slides)}\n")

for slide_idx, slide in enumerate(prs.slides):
    print(f"\n=== Slide {slide_idx} ===")
    for shape_idx, shape in enumerate(slide.shapes):
        print(f"Shape {shape_idx}: {shape.shape_type}")
        if hasattr(shape, "text"):
            print(f"  Text: {shape.text[:100]}")
        if hasattr(shape, "table"):
            try:
                table = shape.table
                print(f"  Table with {len(table.rows)} rows, {len(table.columns)} columns")
                # Print first few cells
                for row_idx in range(min(3, len(table.rows))):
                    row_data = []
                    for col_idx in range(len(table.columns)):
                        cell_text = table.cell(row_idx, col_idx).text
                        row_data.append(cell_text)
                    print(f"    Row {row_idx}: {row_data}")
            except:
                pass
