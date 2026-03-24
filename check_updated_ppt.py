from pptx import Presentation

prs = Presentation('MPL cricket.pptx')

print("All players from updated PowerPoint:\n")
print("Serial | Name")
print("-" * 50)

players_dict = {}
for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text and text not in ["", "Player List"]:
                # Try to parse serial and name
                parts = text.split(maxsplit=1)
                if len(parts) >= 2:
                    try:
                        serial = int(parts[0])
                        name = parts[1]
                        players_dict[serial] = name
                    except:
                        pass

# Sort and print
for serial in sorted(players_dict.keys()):
    print(f"{serial:3d}    | {players_dict[serial]}")

print(f"\nTotal: {len(players_dict)} players")
