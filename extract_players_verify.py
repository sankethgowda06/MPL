import os
from pptx import Presentation
import json

# Path to the PowerPoint file
ppt_path = r'c:\Users\Sanketh\Desktop\MPL\app\MPL cricket.pptx'
images_dir = r'c:\Users\Sanketh\Desktop\MPL\app\player_photos'

# Create images directory if it doesn't exist
os.makedirs(images_dir, exist_ok=True)

# Load the PowerPoint presentation
prs = Presentation(ppt_path)

# Extract player data
players = []

for slide_idx, slide in enumerate(prs.slides):
    print(f"\n=== Slide {slide_idx + 1} ===")
    
    # Extract text from shapes
    texts = []
    image_path = None
    
    for shape in slide.shapes:
        # Extract text
        if hasattr(shape, 'text') and shape.text.strip():
            text = shape.text.strip()
            texts.append(text)
            print(f"Text: {text}")
        
        # Extract images
        if shape.shape_type == 13:  # Picture shape type
            try:
                image = shape.image
                image_bytes = image.blob
                image_filename = f'player_{slide_idx + 1}.{image.ext}'
                image_filepath = os.path.join(images_dir, image_filename)
                
                with open(image_filepath, 'wb') as f:
                    f.write(image_bytes)
                
                image_path = image_filepath
                print(f"Photo saved: {image_filename}")
            except Exception as e:
                print(f"Error extracting image: {e}")
    
    # Parse texts to identify player name and role
    player_name = None
    role = None
    
    # Look for role keywords
    role_keywords = ['all-rounder', 'allrounder', 'batsman', 'bowler', 'keeper', 'wicket keeper', 'spinner']
    
    for text in texts:
        text_lower = text.lower()
        is_role = any(keyword in text_lower for keyword in role_keywords)
        
        if is_role and role is None:
            role = text
        elif not is_role and player_name is None:
            player_name = text
    
    # Add player to list if we have at least a name
    if player_name:
        players.append({
            'serial_number': slide_idx + 1,
            'name': player_name,
            'role': role or 'Unknown',
            'photo_path': image_path
        })
        print(f"✓ Player: {player_name} | Role: {role or 'Unknown'}")

# Save extracted players to JSON for verification
output_file = r'c:\Users\Sanketh\Desktop\MPL\app\extracted_players.json'
with open(output_file, 'w') as f:
    json.dump(players, f, indent=2)

print(f"\n{'='*50}")
print(f"✓ Total players extracted: {len(players)}")
print(f"✓ Data saved to: {output_file}")
print(f"✓ Photos saved to: {images_dir}")
print(f"{'='*50}")

# Display summary
print("\nExtracted Players Summary:")
for player in players:
    print(f"{player['serial_number']}. {player['name']} ({player['role']})")
