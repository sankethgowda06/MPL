from pptx import Presentation

prs = Presentation('player_list.pptx')

print("All players from PowerPoint:\n")
all_text = []
for slide_idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text and text not in ["", "Player List"]:
                all_text.append(text)

# Print all text
for item in all_text:
    print(item)
