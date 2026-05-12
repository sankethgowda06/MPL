import json
import os
import sqlite3
from pathlib import Path

from pptx import Presentation


ROOT_DIR = Path(__file__).resolve().parent
PPT_PATH = ROOT_DIR / "MPL cricket.pptx"
PHOTOS_DIR = ROOT_DIR / "player_photos"
JSON_PATH = ROOT_DIR / "extracted_players.json"
DB_PATH = ROOT_DIR / "instance" / "mpl_league.db"

ROLE_VALUES = {"ALL-ROUNDER", "BATSMAN", "BOWLER", "WICKETKEEPER"}


def parse_player_from_slide(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            value = shape.text.strip()
            if value:
                texts.append(value)

    name = None
    role = None
    for text in texts:
        normalized = text.upper().replace(" ", "")
        if normalized in {"ALL-ROUNDER", "ALLROUNDER"}:
            role = "ALL-ROUNDER"
        elif text.upper() in ROLE_VALUES:
            role = text.upper()
        elif name is None:
            name = text.upper()

    return name, role or "UNKNOWN"


def extract_players_and_photos():
    PHOTOS_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(PPT_PATH))
    players = []

    for idx, slide in enumerate(prs.slides, start=1):
        photo_path = None
        slide_area = prs.slide_width * prs.slide_height
        picture_shapes = []

        for shape in slide.shapes:
            if shape.shape_type == 13:
                shape_area = shape.width * shape.height
                picture_shapes.append((shape_area, shape))

        # Prefer player image over full-slide background image:
        # choose the largest picture that is NOT near full-slide size.
        candidate_shape = None
        non_background = [
            item for item in picture_shapes if item[0] < (slide_area * 0.8)
        ]
        if non_background:
            candidate_shape = max(non_background, key=lambda item: item[0])[1]
        elif picture_shapes:
            candidate_shape = max(picture_shapes, key=lambda item: item[0])[1]

        if candidate_shape is not None:
            image = candidate_shape.image
            extension = image.ext.lower()
            photo_file = PHOTOS_DIR / f"player_{idx}.{extension}"
            with open(photo_file, "wb") as f:
                f.write(image.blob)
            photo_path = str(photo_file)

        name, role = parse_player_from_slide(slide)
        if not name:
            raise ValueError(f"Could not parse player name for slide {idx}")

        players.append(
            {
                "serial_number": idx,
                "name": name,
                "role": role,
                "photo_path": photo_path,
            }
        )

    with open(JSON_PATH, "w", encoding="utf-8") as f:
        json.dump(players, f, indent=2)

    return players


def sync_players_to_db(players):
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cur = conn.cursor()

    cur.execute("SELECT id, serial_number, name, role, photo_path FROM player")
    existing_rows = cur.fetchall()
    by_serial = {row["serial_number"]: row for row in existing_rows}

    added = 0
    updated = 0
    unchanged = 0

    for p in players:
        serial = p["serial_number"]
        row = by_serial.get(serial)

        if row is None:
            cur.execute(
                """
                INSERT INTO player (serial_number, name, role, photo_path, is_available)
                VALUES (?, ?, ?, ?, 1)
                """,
                (serial, p["name"], p["role"], p["photo_path"]),
            )
            added += 1
            continue

        has_change = (
            row["name"] != p["name"]
            or row["role"] != p["role"]
            or (row["photo_path"] or "") != (p["photo_path"] or "")
        )
        if has_change:
            cur.execute(
                """
                UPDATE player
                SET name = ?, role = ?, photo_path = ?
                WHERE id = ?
                """,
                (p["name"], p["role"], p["photo_path"], row["id"]),
            )
            updated += 1
        else:
            unchanged += 1

    conn.commit()
    conn.close()
    return added, updated, unchanged


def main():
    if not PPT_PATH.exists():
        raise FileNotFoundError(f"PPT not found: {PPT_PATH}")
    if not DB_PATH.exists():
        raise FileNotFoundError(f"DB not found: {DB_PATH}")

    players = extract_players_and_photos()
    added, updated, unchanged = sync_players_to_db(players)

    print(f"Extracted players: {len(players)}")
    print(f"Database added: {added}")
    print(f"Database updated: {updated}")
    print(f"Database unchanged: {unchanged}")
    print(f"JSON updated: {JSON_PATH}")
    print(f"Photos updated: {PHOTOS_DIR}")
    print(f"Database updated: {DB_PATH}")


if __name__ == "__main__":
    main()
