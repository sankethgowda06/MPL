import os
from pptx import Presentation
from pptx.util import Inches
import openpyxl
from openpyxl.drawing.image import Image as XLImage
from openpyxl.utils import get_column_letter
import re

# Path to the PowerPoint file
ppt_path = r'c:\Users\Sanketh\Desktop\MPL\app\MPL cricket.pptx'
excel_path = r'c:\Users\Sanketh\Desktop\MPL\app\Player_List.xlsx'
images_dir = r'c:\Users\Sanketh\Desktop\MPL\app\player_photos'

# Create images directory if it doesn't exist
os.makedirs(images_dir, exist_ok=True)

# Load the PowerPoint presentation
prs = Presentation(ppt_path)

# Extract player data
players = []
image_count = 0

for slide_idx, slide in enumerate(prs.slides):
    print(f"\nProcessing Slide {slide_idx + 1}...")
    
    # Extract text from shapes
    texts = []
    image_path = None
    
    for shape_idx, shape in enumerate(slide.shapes):
        # Extract text
        if hasattr(shape, 'text') and shape.text.strip():
            text = shape.text.strip()
            texts.append(text)
            print(f"  Text found: {text}")
        
        # Extract images
        if shape.shape_type == 13:  # Picture shape type
            try:
                image = shape.image
                image_bytes = image.blob
                image_filename = f'player_{slide_idx + 1}.{image.ext}'
                image_filepath = os.path.join(images_dir, image_filename)
                
                with open(image_filepath, 'wb') as f:
                    f.write(image_bytes)
                
                image_path = image_filepath
                print(f"  Image extracted: {image_filename}")
                image_count += 1
            except Exception as e:
                print(f"  Error extracting image: {e}")
    
    # Parse texts to identify player name and role
    player_name = None
    role = None
    
    # Look for role keywords
    role_keywords = ['all-rounder', 'allrounder', 'batsman', 'bowler', 'keeper', 'wicket keeper', 'spinner']
    
    for text in texts:
        text_lower = text.lower()
        # Check if this text is a role
        is_role = any(keyword in text_lower for keyword in role_keywords)
        
        if is_role and role is None:
            role = text
        elif not is_role and player_name is None:
            player_name = text
    
    # Add player to list if we have at least a name
    if player_name:
        serial_no = len(players) + 1
        players.append({
            'Serial Number': serial_no,
            'Player Name': player_name,
            'Role': role or 'Unknown',
            'Photo Path': image_path
        })
        print(f"  Player added: {player_name} - {role or 'Unknown'}")

print(f"\n\nTotal players extracted: {len(players)}")
print(f"Total images extracted: {image_count}")

# Create Excel workbook
wb = openpyxl.Workbook()
ws = wb.active
ws.title = "Players"

# Add headers
headers = ['Serial Number', 'Player Name', 'Role', 'Player Photo']
ws.append(headers)

# Style header row
from openpyxl.styles import Font, PatternFill
header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
header_font = Font(bold=True, color="FFFFFF")

for cell in ws[1]:
    cell.fill = header_fill
    cell.font = header_font

# Add player data
for player in players:
    ws.append([
        player['Serial Number'],
        player['Player Name'],
        player['Role'],
        ''  # Placeholder for image
    ])

# Adjust column widths
ws.column_dimensions['A'].width = 15
ws.column_dimensions['B'].width = 25
ws.column_dimensions['C'].width = 20
ws.column_dimensions['D'].width = 20

# Add images to Excel
for row_idx, player in enumerate(players, start=2):
    if player['Photo Path'] and os.path.exists(player['Photo Path']):
        try:
            img = XLImage(player['Photo Path'])
            img.width = 100
            img.height = 100
            ws.add_image(img, f'D{row_idx}')
            print(f"Image added for {player['Player Name']}")
        except Exception as e:
            print(f"Error adding image for {player['Player Name']}: {e}")

# Adjust row heights for images
for row_idx in range(2, len(players) + 2):
    ws.row_dimensions[row_idx].height = 100

# Save the workbook
wb.save(excel_path)
print(f"\n✓ Excel file created successfully: {excel_path}")
print(f"✓ Player photos saved to: {images_dir}")
